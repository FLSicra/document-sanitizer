from __future__ import annotations
import io
import re
from pathlib import Path
from vault.vault import SanitizeSession


# Pattern matching any token like [ENTITY_TYPE_123]
TOKEN_PATTERN = re.compile(r"\[([A-Z_]+)_(\d+)\]")


def restore_text(text: str, token_map: dict[str, str]) -> str:
    """Replace all tokens in text with their original values."""
    def replace(m: re.Match) -> str:
        token = m.group(0)
        return token_map.get(token, token)
    return TOKEN_PATTERN.sub(replace, text)


def restore_file(sanitized_path: Path, vault_path: Path, password: str, output_path: Path) -> None:
    """
    Load a sanitized document and its vault, then write the restored version to output_path.
    Supports all text-based formats. For binary formats (docx/xlsx/pptx/pdf/odf),
    delegates to the appropriate sanitizer's restore logic.
    """
    token_map = SanitizeSession.load_vault(vault_path, password)
    suffix = sanitized_path.suffix.lower()

    if suffix in (".txt", ".log", ".csv", ".json", ".yaml", ".yml", ".xml",
                  ".ini", ".env", ".toml", ".md"):
        _restore_text_file(sanitized_path, output_path, token_map)
    elif suffix in (".docx", ".dotx"):
        _restore_docx(sanitized_path, output_path, token_map)
    elif suffix in (".xlsx", ".xlsm"):
        _restore_xlsx(sanitized_path, output_path, token_map)
    elif suffix in (".pptx", ".potx"):
        _restore_pptx(sanitized_path, output_path, token_map)
    elif suffix == ".pdf":
        _restore_pdf(sanitized_path, output_path, token_map)
    elif suffix in (".odt", ".ods", ".odp"):
        _restore_odf(sanitized_path, output_path, token_map)
    else:
        # Fallback: treat as text
        _restore_text_file(sanitized_path, output_path, token_map)


def _restore_text_file(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    text = src.read_text(encoding="utf-8", errors="replace")
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(restore_text(text, token_map), encoding="utf-8")


def _restore_docx(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    from docx import Document
    with open(str(src), "rb") as f:
        doc = Document(io.BytesIO(f.read()))
    for para in doc.paragraphs:
        for run in para.runs:
            run.text = restore_text(run.text, token_map)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.text = restore_text(run.text, token_map)
    for section in doc.sections:
        for hf in [
            section.header, section.footer,
            section.even_page_header, section.even_page_footer,
            section.first_page_header, section.first_page_footer,
        ]:
            if hf is not None:
                for para in hf.paragraphs:
                    for run in para.runs:
                        run.text = restore_text(run.text, token_map)
                for table in hf.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.text = restore_text(run.text, token_map)
    dst.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(dst))


def _restore_xlsx(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    from openpyxl import load_workbook
    wb = load_workbook(str(src))
    try:
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        cell.value = restore_text(cell.value, token_map)
        dst.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(dst))
    finally:
        wb.close()


def _restore_pptx(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    from pptx import Presentation
    with open(str(src), "rb") as f:
        prs = Presentation(io.BytesIO(f.read()))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = restore_text(run.text, token_map)
    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))


def _restore_pdf(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    """
    PDF restoration is limited: black redaction boxes cannot be un-drawn.
    We can only restore if the PDF used text overlay (not pixel redaction).
    For now, raise an informative error.
    """
    raise NotImplementedError(
        "PDF visual redaction (black boxes) cannot be reversed. "
        "Keep the original PDF and use the vault only to verify what was redacted."
    )


def _restore_odf(src: Path, dst: Path, token_map: dict[str, str]) -> None:
    from odf.opendocument import load
    doc = load(str(src))

    def walk(node):
        if hasattr(node, 'childNodes'):
            for child in node.childNodes:
                if hasattr(child, 'data') and isinstance(child.data, str):
                    child.data = restore_text(child.data, token_map)
                else:
                    walk(child)

    if hasattr(doc, 'text'):
        walk(doc.text)
    dst.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(dst))
