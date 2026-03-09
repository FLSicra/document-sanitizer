import io
import re
from pathlib import Path
from sanitizers.base import (
    Detection, SanitizeResult, Sanitizer, dedup_detections,
    extract_company_roots, find_company_root_hits,
    replace_detections_in_text,
)
from detectors.engine import analyze_text
from utils.streaming import check_zip_bomb


def _clear_para_content_docx(para) -> None:
    """Remove all content children from a DOCX paragraph, preserving paragraph properties."""
    from docx.oxml.ns import qn
    pPr = para._p.find(qn('w:pPr'))
    for child in list(para._p):
        para._p.remove(child)
    if pPr is not None:
        para._p.insert(0, pPr)


def _clear_para_content_pptx(para) -> None:
    """Remove all content children from a PPTX paragraph, preserving paragraph properties."""
    from pptx.oxml.ns import qn
    pPr = para._p.find(qn('a:pPr'))
    for child in list(para._p):
        para._p.remove(child)
    if pPr is not None:
        para._p.insert(0, pPr)


def _replace_in_text(text: str, detections: list[Detection], session) -> str:
    """Thin wrapper around the shared replacement helper."""
    return replace_detections_in_text(text, detections, session)


def _replace_in_runs(para, detections: list[Detection], session) -> None:
    """Replace detections within individual runs to preserve per-run formatting."""
    to_redact = dedup_detections([d for d in detections if d.redact])
    if not to_redact:
        return
    runs = list(para.runs)
    if not runs:
        return

    # Compute each run's start offset within para.text
    run_starts: list[int] = []
    pos = 0
    for run in runs:
        run_starts.append(pos)
        pos += len(run.text)

    # Process in reverse order so earlier offsets stay valid
    for d in sorted(to_redact, key=lambda d: d.start, reverse=True):
        token = session.get_or_create_token(d) if session else "[REDACTED]"
        d.token = token

        # Find all runs overlapping [d.start, d.end)
        affected: list[tuple[int, int, object]] = []
        for i, run in enumerate(runs):
            rs = run_starts[i]
            run_end = rs + len(run.text)
            if rs < d.end and run_end > d.start:
                affected.append((i, rs, run))

        if not affected:
            continue

        if len(affected) == 1:
            _, rs, run = affected[0]
            local_start = d.start - rs
            local_end = d.end - rs
            run.text = run.text[:local_start] + token + run.text[local_end:]
        else:
            # Detection spans multiple runs: put token in first, clear middle, trim last
            _, rs0, r0 = affected[0]
            r0.text = r0.text[:d.start - rs0] + token
            for _, _, r in affected[1:-1]:
                r.text = ""
            _, rs_last, r_last = affected[-1]
            r_last.text = r_last.text[d.end - rs_last:]


def _collect_detections_from_text(
    text: str,
    context: str,
    custom_terms: tuple[str, ...] = (),
    enabled_entities: frozenset[str] | None = None,
) -> list[Detection]:
    if not text.strip():
        return []
    results = analyze_text(text, custom_terms, enabled_entities)
    out = []
    for r in results:
        out.append(Detection(
            entity_type=r.entity_type,
            original_value=text[r.start:r.end],
            start=r.start,
            end=r.end,
            score=r.score,
            page_or_line=context,
        ))
    return dedup_detections(out)


class OfficeSanitizer(Sanitizer):
    def _get_handler(self):
        suffix = self.path.suffix.lower()
        if suffix in (".docx", ".dotx"):
            return _DocxHandler(self.path)
        elif suffix in (".xlsx", ".xlsm"):
            return _XlsxHandler(self.path)
        elif suffix in (".pptx", ".potx"):
            return _PptxHandler(self.path)
        raise ValueError(f"Unsupported office format: {suffix}")

    def detect(
        self,
        custom_terms: tuple[str, ...] = (),
        enabled_entities: frozenset[str] | None = None,
        progress_callback=None,
    ) -> list[Detection]:
        check_zip_bomb(self.path)
        return self._get_handler().detect(custom_terms, enabled_entities)

    def sanitize(self, detections: list[Detection], output_path: Path, session) -> SanitizeResult:
        check_zip_bomb(self.path)
        try:
            result = self._get_handler().sanitize(detections, output_path, session)
            result.source_path = self.path
            return result
        except Exception as e:
            return SanitizeResult(source_path=self.path, error=str(e))


class _DocxHandler:
    def __init__(self, path: Path):
        self.path = path

    @staticmethod
    def _iter_paragraphs(doc):
        """Yield (paragraph, context_key) for body, tables, and all headers/footers."""
        for i, para in enumerate(doc.paragraphs):
            yield para, f"para {i+1}"
        for t_idx, table in enumerate(doc.tables):
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for p_idx, para in enumerate(cell.paragraphs):
                        yield para, f"table {t_idx} cell {r_idx},{c_idx} para {p_idx}"
        for section in doc.sections:
            for hf_name, hf in [("hdr", section.header), ("ftr", section.footer),
                                 ("ehdr", section.even_page_header), ("eftr", section.even_page_footer),
                                 ("fhdr", section.first_page_header), ("fftr", section.first_page_footer)]:
                if hf is not None:
                    for i, para in enumerate(hf.paragraphs):
                        yield para, f"{hf_name} para {i}"

    def detect(
        self,
        custom_terms: tuple[str, ...] = (),
        enabled_entities: frozenset[str] | None = None,
        progress_callback=None,
    ) -> list[Detection]:
        from docx import Document
        with open(str(self.path), "rb") as f:
            doc = Document(io.BytesIO(f.read()))
        paragraphs = [(para.text, context) for para, context in self._iter_paragraphs(doc)]
        detections = []
        for text, context in paragraphs:
            detections.extend(_collect_detections_from_text(text, context, custom_terms, enabled_entities))
        roots = extract_company_roots(detections)
        if roots:
            for text, context in paragraphs:
                detections.extend(find_company_root_hits(text, roots, context))
            detections = dedup_detections(detections)
        return detections

    def sanitize(self, detections: list[Detection], output_path: Path, session) -> SanitizeResult:
        from docx import Document
        with open(str(self.path), "rb") as f:
            doc = Document(io.BytesIO(f.read()))
        if session is not None:
            session.initialize_from_content(
                [para.text for para, _ in self._iter_paragraphs(doc)]
            )
        detection_map: dict[str, list[Detection]] = {}
        for d in detections:
            detection_map.setdefault(d.page_or_line, []).append(d)

        for para, context in self._iter_paragraphs(doc):
            para_detections = detection_map.get(context, [])
            if para_detections:
                if para.runs:
                    _replace_in_runs(para, para_detections, session)
                else:
                    # Paragraph has no runs — fall back to full-text replacement
                    new_text = _replace_in_text(para.text, para_detections, session)
                    if new_text != para.text:
                        _clear_para_content_docx(para)
                        para.add_run(new_text)

        layout_warnings = [
            f"Token longer than original in '{d.page_or_line}' — table layout may shift."
            for d in detections
            if d.redact and d.token and len(d.token) - len(d.original_value) > 10
        ]
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        result = SanitizeResult(source_path=self.path, output_path=output_path, detections=detections)
        result.warnings = layout_warnings
        return result


class _XlsxHandler:
    def __init__(self, path: Path):
        self.path = path

    def detect(
        self,
        custom_terms: tuple[str, ...] = (),
        enabled_entities: frozenset[str] | None = None,
        progress_callback=None,
    ) -> list[Detection]:
        from openpyxl import load_workbook
        cells = []
        formulas = []
        # Single load with data_only=False: literal string cells keep their
        # values, and formula cells expose raw formula text for inspection.
        wb = load_workbook(str(self.path), read_only=True, data_only=False)
        try:
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if not (cell.value and isinstance(cell.value, str)):
                            continue
                        coord = f"{sheet.title}!{cell.coordinate}"
                        if cell.value.startswith("="):
                            formulas.append((cell.value, coord))
                        else:
                            cells.append((cell.value, coord))
        finally:
            wb.close()
        detections = []
        for text, context in cells:
            detections.extend(_collect_detections_from_text(text, context, custom_terms, enabled_entities))
        roots = extract_company_roots(detections)
        if roots:
            for text, context in cells:
                detections.extend(find_company_root_hits(text, roots, context))
            detections = dedup_detections(detections)
        # Scan for dangerous Excel formulas that make external requests.
        for v, coord in formulas:
            vupper = v.upper()
            if (
                vupper.startswith(("=IMAGE(", "=WEBSERVICE(", "=FILTERXML(", "=IMPORTDATA("))
                or re.search(r"=.*https?://", v, re.IGNORECASE)
            ):
                detections.append(Detection(
                    entity_type="DANGEROUS_FORMULA",
                    original_value=v,
                    start=0,
                    end=len(v),
                    score=1.0,
                    page_or_line=coord,
                    redact=True,
                ))
        return detections

    def sanitize(self, detections: list[Detection], output_path: Path, session) -> SanitizeResult:
        from openpyxl import load_workbook
        wb = load_workbook(str(self.path))
        try:
            if session is not None:
                all_texts = [
                    str(cell.value)
                    for sheet in wb.worksheets
                    for row in sheet.iter_rows()
                    for cell in row
                    if cell.value and isinstance(cell.value, str)
                ]
                session.initialize_from_content(all_texts)
            detection_map: dict[str, list[Detection]] = {}
            for d in detections:
                detection_map.setdefault(d.page_or_line, []).append(d)

            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            key = f"{sheet.title}!{cell.coordinate}"
                            cell_detections = detection_map.get(key, [])
                            if cell_detections:
                                cell.value = _replace_in_text(cell.value, cell_detections, session)
                                # Expand column width to fit the (possibly longer) token
                                col_letter = cell.column_letter
                                current = sheet.column_dimensions[col_letter].width or 8
                                needed = len(cell.value) + 2
                                if needed > current:
                                    sheet.column_dimensions[col_letter].width = needed

            output_path.parent.mkdir(parents=True, exist_ok=True)
            wb.save(str(output_path))
        finally:
            wb.close()
        return SanitizeResult(source_path=self.path, output_path=output_path, detections=detections)


class _PptxHandler:
    def __init__(self, path: Path):
        self.path = path

    def detect(
        self,
        custom_terms: tuple[str, ...] = (),
        enabled_entities: frozenset[str] | None = None,
        progress_callback=None,
    ) -> list[Detection]:
        from pptx import Presentation
        with open(str(self.path), "rb") as f:
            prs = Presentation(io.BytesIO(f.read()))
        paragraphs = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    try:
                        for para_idx, para in enumerate(shape.text_frame.paragraphs):
                            paragraphs.append((para.text, f"slide {slide_idx} para {para_idx+1}"))
                    except (AttributeError, KeyError):
                        continue
        detections = []
        for text, context in paragraphs:
            detections.extend(_collect_detections_from_text(text, context, custom_terms, enabled_entities))
        roots = extract_company_roots(detections)
        if roots:
            for text, context in paragraphs:
                detections.extend(find_company_root_hits(text, roots, context))
            detections = dedup_detections(detections)
        return detections

    def sanitize(self, detections: list[Detection], output_path: Path, session) -> SanitizeResult:
        from pptx import Presentation
        with open(str(self.path), "rb") as f:
            prs = Presentation(io.BytesIO(f.read()))
        if session is not None:
            all_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        try:
                            for para in shape.text_frame.paragraphs:
                                if para.text:
                                    all_texts.append(para.text)
                        except (AttributeError, KeyError):
                            pass
            session.initialize_from_content(all_texts)
        detection_map: dict[str, list[Detection]] = {}
        for d in detections:
            detection_map.setdefault(d.page_or_line, []).append(d)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        key = f"slide {slide_idx} para {para_idx+1}"
                        para_detections = detection_map.get(key, [])
                        if para_detections:
                            if para.runs:
                                _replace_in_runs(para, para_detections, session)
                            else:
                                new_text = _replace_in_text(para.text, para_detections, session)
                                if new_text != para.text:
                                    _clear_para_content_pptx(para)
                                    run = para.add_run()
                                    run.text = new_text

        layout_warnings = [
            f"Token longer than original in '{d.page_or_line}' — table layout may shift."
            for d in detections
            if d.redact and d.token and len(d.token) - len(d.original_value) > 10
        ]
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        result = SanitizeResult(source_path=self.path, output_path=output_path, detections=detections)
        result.warnings = layout_warnings
        return result
